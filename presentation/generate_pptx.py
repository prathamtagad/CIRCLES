"""
🎯 Circles Hackathon Presentation Generator
Generates a stunning PowerPoint presentation for SGSITS Hackathon 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor
import os

# Color palette matching Circles app
COLORS = {
    'bg_dark': RGBColor(3, 7, 18),           # #030712
    'primary': RGBColor(139, 92, 246),       # #8B5CF6 Violet
    'accent': RGBColor(6, 182, 212),         # #06B6D4 Cyan
    'inner': RGBColor(34, 211, 238),         # #22D3EE
    'middle': RGBColor(129, 140, 248),       # #818CF8
    'outer': RGBColor(148, 163, 184),        # #94A3B8
    'danger': RGBColor(244, 63, 94),         # #F43F5E
    'success': RGBColor(16, 185, 129),       # #10B981
    'text': RGBColor(248, 250, 252),         # #F8FAFC
    'text_muted': RGBColor(148, 163, 184),   # #94A3B8
    'surface': RGBColor(15, 23, 42),         # #0F172A
    'pink': RGBColor(236, 72, 153),          # #EC4899
    'amber': RGBColor(251, 191, 36),         # #FBB424
}

def set_slide_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    """Add a text box with styling"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS['text']
    p.alignment = align
    return txBox

def add_shape_with_text(slide, shape_type, left, top, width, height, text="", fill_color=None, line_color=None, font_size=14, font_color=None):
    """Add a shape with optional text"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color or COLORS['text']
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.bold = True
    
    return shape

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    
    # ========== SLIDE 1: Title ==========
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLORS['bg_dark'])
    
    # Logo circle
    logo = add_shape_with_text(slide1, MSO_SHAPE.OVAL, 5.9, 0.8, 1.5, 1.5, "", COLORS['primary'])
    
    # Inner circle
    inner_logo = add_shape_with_text(slide1, MSO_SHAPE.OVAL, 6.15, 1.05, 1.0, 1.0, "", COLORS['text'])
    
    # Title
    add_text_box(slide1, 0, 2.5, 13.333, 1, "Circles", 72, True, COLORS['text'], PP_ALIGN.CENTER)
    
    # Subtitle
    add_text_box(slide1, 0, 3.6, 13.333, 0.6, "Redefining Social Connection with Radical Privacy", 28, False, COLORS['primary'], PP_ALIGN.CENTER)
    
    # Tagline
    add_text_box(slide1, 0, 4.3, 13.333, 0.5, 'Share what you want, ONLY with who you want.', 20, False, COLORS['text_muted'], PP_ALIGN.CENTER)
    
    # Hackathon badge
    badge = add_shape_with_text(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, 4.5, 5.2, 4.3, 0.6, "🏆 SGSITS Hackathon 2026", COLORS['surface'], COLORS['primary'], 16)
    
    # Slide number
    add_text_box(slide1, 12.5, 7, 0.5, 0.3, "01", 12, False, COLORS['text_muted'], PP_ALIGN.RIGHT)
    
    # ========== SLIDE 2: The Problem ==========
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLORS['bg_dark'])
    
    # Header
    add_text_box(slide2, 0.5, 0.4, 8, 0.8, "The Privacy Paradox 🚨", 44, True, COLORS['text'], PP_ALIGN.LEFT)
    
    # Underline
    add_shape_with_text(slide2, MSO_SHAPE.RECTANGLE, 0.5, 1.15, 5, 0.08, "", COLORS['primary'])
    
    # Problem cards
    problems = [
        ("🔓", "All-or-Nothing", "Current platforms force you to share with everyone or be completely invisible.", 0.5),
        ("😰", "Context Collapse", "Your boss, grandma, and best friend see the same content.", 4.7),
        ("📊", "Data Exploitation", "Users have lost control over their digital footprint.", 8.9),
    ]
    
    for icon, title, desc, left in problems:
        # Card background
        card = add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE, left, 1.8, 3.9, 4.5, "", COLORS['surface'], COLORS['danger'])
        
        # Icon
        add_text_box(slide2, left + 0.2, 2.1, 1, 0.8, icon, 48, False, COLORS['danger'], PP_ALIGN.LEFT)
        
        # Title
        add_text_box(slide2, left + 0.2, 3.0, 3.5, 0.5, title, 22, True, COLORS['danger'], PP_ALIGN.LEFT)
        
        # Description
        add_text_box(slide2, left + 0.2, 3.6, 3.5, 1.5, desc, 14, False, COLORS['text_muted'], PP_ALIGN.LEFT)
    
    add_text_box(slide2, 12.5, 7, 0.5, 0.3, "02", 12, False, COLORS['text_muted'], PP_ALIGN.RIGHT)
    
    # ========== SLIDE 3: The Solution ==========
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLORS['bg_dark'])
    
    # Header
    add_text_box(slide3, 0.5, 0.4, 8, 0.8, "Introducing Circles 💫", 44, True, COLORS['text'], PP_ALIGN.LEFT)
    add_shape_with_text(slide3, MSO_SHAPE.RECTANGLE, 0.5, 1.15, 5, 0.08, "", COLORS['primary'])
    
    # Concentric circles visualization
    outer = add_shape_with_text(slide3, MSO_SHAPE.OVAL, 1, 1.8, 4.5, 4.5, "", None, COLORS['outer'])
    middle = add_shape_with_text(slide3, MSO_SHAPE.OVAL, 1.75, 2.55, 3, 3, "", None, COLORS['middle'])
    inner = add_shape_with_text(slide3, MSO_SHAPE.OVAL, 2.5, 3.3, 1.5, 1.5, "", None, COLORS['inner'])
    core = add_shape_with_text(slide3, MSO_SHAPE.OVAL, 2.9, 3.7, 0.7, 0.7, "", COLORS['primary'])
    
    # Labels on the right
    circles_info = [
        ("🔒", "Inner Circle", "Encrypted • Closest friends • Full trust", COLORS['inner'], 2.0),
        ("🛡️", "Trusted Circle", "Friends & family • Shared memories", COLORS['middle'], 3.5),
        ("👥", "Extended Circle", "Acquaintances • Professional contacts", COLORS['outer'], 5.0),
    ]
    
    for icon, title, desc, color, top in circles_info:
        # Card
        add_shape_with_text(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, 6, top, 6.8, 1.2, "", COLORS['surface'], None)
        
        # Color indicator
        add_shape_with_text(slide3, MSO_SHAPE.OVAL, 6.3, top + 0.4, 0.4, 0.4, "", color)
        
        # Text
        add_text_box(slide3, 7, top + 0.2, 5.5, 0.4, f"{icon} {title}", 18, True, COLORS['text'], PP_ALIGN.LEFT)
        add_text_box(slide3, 7, top + 0.65, 5.5, 0.4, desc, 12, False, COLORS['text_muted'], PP_ALIGN.LEFT)
    
    add_text_box(slide3, 12.5, 7, 0.5, 0.3, "03", 12, False, COLORS['text_muted'], PP_ALIGN.RIGHT)
    
    # ========== SLIDE 4: Killer Features ==========
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLORS['bg_dark'])
    
    add_text_box(slide4, 0.5, 0.4, 8, 0.8, "Killer Features 🚀", 44, True, COLORS['text'], PP_ALIGN.LEFT)
    add_shape_with_text(slide4, MSO_SHAPE.RECTANGLE, 0.5, 1.15, 4, 0.08, "", COLORS['primary'])
    
    features = [
        ("📊", "Privacy Score", "Gamified metric showing how protected your digital footprint is", COLORS['primary'], 0.5, 1.8),
        ("⚡", "Lite Mode", "90% bandwidth reduction. Perfect for developing regions", COLORS['accent'], 6.9, 1.8),
        ("🗑️", "Right to be Forgotten", "One-click delete or export. True data ownership", COLORS['danger'], 0.5, 4.4),
        ("👁️", "Transparency Dashboard", "Real-time stats on data requests. Nothing to hide", COLORS['success'], 6.9, 4.4),
    ]
    
    for icon, title, desc, color, left, top in features:
        # Card
        add_shape_with_text(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, 5.9, 2.3, "", COLORS['surface'], None)
        
        # Icon box
        add_shape_with_text(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, left + 0.3, top + 0.3, 0.9, 0.9, icon, color, None, 28)
        
        # Title
        add_text_box(slide4, left + 0.3, top + 1.3, 5.3, 0.4, title, 20, True, COLORS['text'], PP_ALIGN.LEFT)
        
        # Description
        add_text_box(slide4, left + 0.3, top + 1.7, 5.3, 0.8, desc, 12, False, COLORS['text_muted'], PP_ALIGN.LEFT)
    
    add_text_box(slide4, 12.5, 7, 0.5, 0.3, "04", 12, False, COLORS['text_muted'], PP_ALIGN.RIGHT)
    
    # ========== SLIDE 5: Live Demo ==========
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLORS['bg_dark'])
    
    add_text_box(slide5, 0.5, 0.4, 8, 0.8, "Live Demo 🎬", 44, True, COLORS['text'], PP_ALIGN.LEFT)
    add_shape_with_text(slide5, MSO_SHAPE.RECTANGLE, 0.5, 1.15, 3, 0.08, "", COLORS['primary'])
    
    add_text_box(slide5, 0.5, 1.4, 10, 0.5, 'See Circles in action with our "Neon Glass" aesthetic', 18, False, COLORS['text_muted'], PP_ALIGN.LEFT)
    
    # Demo steps
    steps = [
        ("1", "🔑", "Login", "Google Auth"),
        ("2", "✍️", "Create Post", "Inner Circle only"),
        ("3", "⚡", "Toggle Lite", "90% savings"),
        ("4", "🛡️", "Privacy Score", "Check dashboard"),
    ]
    
    for i, (num, icon, title, desc) in enumerate(steps):
        left = 0.5 + (i * 3.2)
        
        # Card
        add_shape_with_text(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, left, 2.2, 2.9, 2.2, "", COLORS['surface'], None)
        
        # Number badge
        add_shape_with_text(slide5, MSO_SHAPE.OVAL, left + 1.1, 1.95, 0.5, 0.5, num, COLORS['primary'], None, 14)
        
        # Icon
        add_text_box(slide5, left, 2.5, 2.9, 0.8, icon, 36, False, COLORS['text'], PP_ALIGN.CENTER)
        
        # Title
        add_text_box(slide5, left, 3.3, 2.9, 0.4, title, 16, True, COLORS['text'], PP_ALIGN.CENTER)
        
        # Desc
        add_text_box(slide5, left, 3.7, 2.9, 0.4, desc, 11, False, COLORS['text_muted'], PP_ALIGN.CENTER)
    
    # Stats
    stats = [
        ("100+", "Privacy Score"),
        ("90%", "Bandwidth Saved"),
        ("0", "Data Sold"),
    ]
    
    for i, (value, label) in enumerate(stats):
        left = 2.5 + (i * 3)
        add_text_box(slide5, left, 4.8, 2.5, 0.8, value, 48, True, COLORS['primary'], PP_ALIGN.CENTER)
        add_text_box(slide5, left, 5.6, 2.5, 0.4, label, 14, False, COLORS['text_muted'], PP_ALIGN.CENTER)
    
    add_text_box(slide5, 12.5, 7, 0.5, 0.3, "05", 12, False, COLORS['text_muted'], PP_ALIGN.RIGHT)
    
    # ========== SLIDE 6: Tech Stack ==========
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLORS['bg_dark'])
    
    add_text_box(slide6, 0.5, 0.4, 8, 0.8, "Tech Stack ⚙️", 44, True, COLORS['text'], PP_ALIGN.LEFT)
    add_shape_with_text(slide6, MSO_SHAPE.RECTANGLE, 0.5, 1.15, 3, 0.08, "", COLORS['primary'])
    add_text_box(slide6, 0.5, 1.4, 10, 0.5, "Built for Speed & Security", 18, False, COLORS['text_muted'], PP_ALIGN.LEFT)
    
    tech = [
        ("⚛️", "React + Vite", "Lightning-fast Frontend"),
        ("🎨", "TailwindCSS", "Custom Glassmorphism"),
        ("🔥", "Firebase", "Auth & Realtime DB"),
        ("🔐", "Firestore Rules", "Security at Core"),
    ]
    
    for i, (icon, name, role) in enumerate(tech):
        left = 0.5 + (i * 3.2)
        
        # Card
        add_shape_with_text(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, left, 2.2, 2.9, 3.5, "", COLORS['surface'], None)
        
        # Icon
        add_text_box(slide6, left, 2.5, 2.9, 1, icon, 48, False, COLORS['text'], PP_ALIGN.CENTER)
        
        # Name
        add_text_box(slide6, left, 3.6, 2.9, 0.5, name, 18, True, COLORS['text'], PP_ALIGN.CENTER)
        
        # Role
        add_text_box(slide6, left, 4.1, 2.9, 0.5, role, 12, False, COLORS['text_muted'], PP_ALIGN.CENTER)
    
    add_text_box(slide6, 12.5, 7, 0.5, 0.3, "06", 12, False, COLORS['text_muted'], PP_ALIGN.RIGHT)
    
    # ========== SLIDE 7: Team ==========
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLORS['bg_dark'])
    
    add_text_box(slide7, 0.5, 0.4, 8, 0.8, "Meet the Minds 👥", 44, True, COLORS['text'], PP_ALIGN.LEFT)
    add_shape_with_text(slide7, MSO_SHAPE.RECTANGLE, 0.5, 1.15, 4, 0.08, "", COLORS['primary'])
    
    team = [
        ("🧠", "Pratham Tagad", "Lead Architect & Full Stack", "Vision & Architecture", COLORS['amber'], True),
        ("🎨", "Karan Sachdev", "Frontend & UI/UX", "Pixel-perfect interfaces", COLORS['accent'], False),
        ("🔐", "Madhav Mandhanya", "Backend & Security", "Data privacy & security", COLORS['success'], False),
    ]
    
    for i, (icon, name, role, bio, color, is_lead) in enumerate(team):
        left = 0.5 + (i * 4.2)
        
        # Card with special styling for lead
        border = color if is_lead else None
        add_shape_with_text(slide7, MSO_SHAPE.ROUNDED_RECTANGLE, left, 1.7, 3.9, 4.8, "", COLORS['surface'], border)
        
        # Lead badge
        if is_lead:
            add_shape_with_text(slide7, MSO_SHAPE.ROUNDED_RECTANGLE, left + 1.0, 1.45, 1.9, 0.45, "🏆 TEAM LEAD", COLORS['amber'], None, 10)
        
        # Avatar
        add_shape_with_text(slide7, MSO_SHAPE.ROUNDED_RECTANGLE, left + 1.2, 2.2, 1.5, 1.5, icon, color, None, 36)
        
        # Name
        add_text_box(slide7, left + 0.2, 3.9, 3.5, 0.5, name, 20, True, COLORS['text'], PP_ALIGN.CENTER)
        
        # Role
        add_text_box(slide7, left + 0.2, 4.4, 3.5, 0.5, role, 13, True, COLORS['primary'], PP_ALIGN.CENTER)
        
        # Bio
        add_text_box(slide7, left + 0.2, 4.9, 3.5, 0.8, bio, 11, False, COLORS['text_muted'], PP_ALIGN.CENTER)
    
    add_text_box(slide7, 12.5, 7, 0.5, 0.3, "07", 12, False, COLORS['text_muted'], PP_ALIGN.RIGHT)
    
    # ========== SLIDE 8: Conclusion ==========
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, COLORS['bg_dark'])
    
    add_text_box(slide8, 0, 1.5, 13.333, 1.2, "The Future is Private", 64, True, COLORS['text'], PP_ALIGN.CENTER)
    add_text_box(slide8, 0, 2.8, 13.333, 0.6, "Privacy isn't a feature — it's the foundation.", 24, False, COLORS['text_muted'], PP_ALIGN.CENTER)
    
    # CTA Button
    add_shape_with_text(slide8, MSO_SHAPE.ROUNDED_RECTANGLE, 4.5, 4, 4.3, 0.9, "🗳️ Vote for Circles!", COLORS['primary'], None, 24)
    
    # Footer
    add_text_box(slide8, 0, 5.5, 13.333, 0.5, "Built with ❤️ at SGSITS Hackathon 2026", 16, False, COLORS['text_muted'], PP_ALIGN.CENTER)
    
    add_text_box(slide8, 12.5, 7, 0.5, 0.3, "08", 12, False, COLORS['text_muted'], PP_ALIGN.RIGHT)
    
    # ========== SLIDE 9: Thank You ==========
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, COLORS['bg_dark'])
    
    # Logo
    add_shape_with_text(slide9, MSO_SHAPE.OVAL, 5.9, 1.2, 1.5, 1.5, "", COLORS['primary'])
    add_shape_with_text(slide9, MSO_SHAPE.OVAL, 6.15, 1.45, 1.0, 1.0, "", COLORS['text'])
    
    add_text_box(slide9, 0, 3.0, 13.333, 1, "Thank You!", 60, True, COLORS['text'], PP_ALIGN.CENTER)
    add_text_box(slide9, 0, 4.0, 13.333, 0.6, "Questions? Let's Connect!", 24, False, COLORS['primary'], PP_ALIGN.CENTER)
    add_text_box(slide9, 0, 5.0, 13.333, 0.5, "Circles — Where trust meets technology", 18, False, COLORS['text_muted'], PP_ALIGN.CENTER)
    
    add_text_box(slide9, 12.5, 7, 0.5, 0.3, "09", 12, False, COLORS['text_muted'], PP_ALIGN.RIGHT)
    
    return prs

def main():
    print("Generating Circles Hackathon Presentation...")
    
    # Create output directory
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "Circles_SGSITS_Hackathon_2026.pptx")
    
    # Generate presentation
    prs = create_presentation()
    
    # Save
    prs.save(output_path)
    print("Presentation saved to: " + output_path)
    print("Go win that hackathon!")

if __name__ == "__main__":
    main()

